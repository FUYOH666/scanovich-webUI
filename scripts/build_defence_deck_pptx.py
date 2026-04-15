#!/usr/bin/env python3
"""Build docs/submission/GPTHub_defence_10slides.pptx from docs/submission/SLIDES_10_RU.md.

Requires: python-pptx, Pillow (optional, for image aspect scaling).

Run from repo root:
  uv run --with python-pptx --with pillow python scripts/build_defence_deck_pptx.py

Export to PDF for the hackathon form: Keynote / PowerPoint / LibreOffice
(File → Export as PDF). Typical output size is well under 20 MiB.
"""

from __future__ import annotations

import argparse
import logging
import re
import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger("build_defence_deck_pptx")

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_MD = ROOT / "docs" / "submission" / "SLIDES_10_RU.md"
DEFAULT_OUT = ROOT / "docs" / "submission" / "GPTHub_defence_10slides.pptx"
SUBMISSION_DIR = ROOT / "docs" / "submission"

SLIDE_HEADER_RE = re.compile(r"^##\s+Слайд\s+(\d+)\s+—\s+(.+?)\s*$")


@dataclass
class SlideSpec:
    number: int
    title: str
    image: Path | None
    bullets: list[str]


def parse_slides_md(text: str, base_dir: Path) -> list[SlideSpec]:
    """Parse SLIDES_10_RU.md blocks: ## Слайд N — title, optional @image, bullets."""
    specs: list[SlideSpec] = []
    current_num: int | None = None
    current_title: str | None = None
    current_image: Path | None = None
    current_bullets: list[str] = []

    def flush() -> None:
        nonlocal current_num, current_title, current_image, current_bullets
        if current_num is None:
            return
        if current_title is None:
            raise ValueError(f"Slide {current_num}: missing title")
        specs.append(
            SlideSpec(
                number=current_num,
                title=current_title,
                image=current_image,
                bullets=list(current_bullets),
            )
        )
        current_num = None
        current_title = None
        current_image = None
        current_bullets = []

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        m = SLIDE_HEADER_RE.match(line)
        if m:
            flush()
            current_num = int(m.group(1))
            current_title = m.group(2).strip()
            current_image = None
            current_bullets = []
            continue
        if current_num is None:
            continue
        stripped = line.strip()
        if not stripped or stripped == "---":
            continue
        if stripped.startswith("@image "):
            name = stripped.split(maxsplit=1)[1].strip()
            current_image = (base_dir / name).resolve()
            continue
        if stripped.startswith("- "):
            current_bullets.append(stripped[2:].strip())
            continue
        # allow continuation without bullet (rare) — treat as extra bullet
        if stripped.startswith("*"):
            continue
        logger.warning("ignored line in slide %s: %r", current_num, raw_line[:120])

    flush()
    specs.sort(key=lambda s: s.number)
    if len(specs) != 10:
        raise SystemExit(f"expected 10 slides, got {len(specs)} — check markdown headers")
    for i, s in enumerate(specs, start=1):
        if s.number != i:
            raise SystemExit(f"slides must be numbered 1..10 in order; got slide number {s.number} at position {i}")
        if s.image is not None and not s.image.is_file():
            raise SystemExit(f"slide {s.number}: image not found: {s.image}")
        if not s.bullets and s.image is None:
            raise SystemExit(f"slide {s.number}: add bullets or @image")
    return specs


def _layout_index(prs: Presentation, *candidates: int) -> int:
    for idx in candidates:
        if 0 <= idx < len(prs.slide_layouts):
            return idx
    return 1


def _picture_size(path: Path, max_w_in: float, max_h_in: float) -> tuple[float, float]:
    """Return (width, height) in inches preserving aspect ratio."""
    try:
        from PIL import Image as PILImage

        with PILImage.open(path) as im:
            w_px, h_px = im.size
    except Exception:
        return max_w_in, max_h_in * 0.75
    aspect = w_px / max(float(h_px), 1.0)
    w_in = max_w_in
    h_in = w_in / aspect
    if h_in > max_h_in:
        h_in = max_h_in
        w_in = h_in * aspect
    return w_in, h_in


def build_pptx(specs: list[SlideSpec], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = int(Inches(13.333))
    prs.slide_height = int(Inches(7.5))

    layout_title_content = _layout_index(prs, 1, 3)
    layout_blank = _layout_index(prs, 6, 5)

    for spec in specs:
        if spec.image is not None:
            slide = prs.slides.add_slide(prs.slide_layouts[layout_blank])
            title_box = slide.shapes.add_textbox(
                Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.75)
            )
            tf = title_box.text_frame
            tf.text = spec.title
            p0 = tf.paragraphs[0]
            p0.font.size = Pt(26)
            p0.font.bold = True

            max_w, max_h = 12.35, 5.95
            wi, hi = _picture_size(spec.image, max_w, max_h)
            left = (13.333 - wi) / 2
            top = 1.05 + (5.95 - hi) / 2 * 0.15  # nudge down slightly
            slide.shapes.add_picture(str(spec.image), Inches(left), Inches(top), Inches(wi), Inches(hi))

            if spec.bullets:
                cap = slide.shapes.add_textbox(
                    Inches(0.45), Inches(6.85), Inches(12.4), Inches(0.55)
                )
                ctf = cap.text_frame
                ctf.text = spec.bullets[0][:500]
                ctf.paragraphs[0].font.size = Pt(11)
                if len(spec.bullets) > 1:
                    for extra in spec.bullets[1:3]:
                        pp = ctf.add_paragraph()
                        pp.text = extra[:300]
                        pp.font.size = Pt(11)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[layout_title_content])
            title = slide.shapes.title
            title.text = spec.title
            for p in title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(32)
                    r.font.bold = True

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(spec.bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(17) if len(spec.bullets) > 5 else Pt(19)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    logger.info("wrote %s (%s slides)", out_path.relative_to(ROOT), len(prs.slides))


def main(argv: list[str] | None = None) -> int:
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s %(levelname)s %(name)s %(message)s",
    )
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--md", type=Path, default=DEFAULT_MD, help="Path to SLIDES_10_RU.md")
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT, help="Output .pptx path")
    parser.add_argument(
        "--submission-dir",
        type=Path,
        default=SUBMISSION_DIR,
        help="Directory for resolving @image paths",
    )
    args = parser.parse_args(argv)
    md_path = args.md.resolve()
    if not md_path.is_file():
        logger.error("markdown not found: %s", md_path)
        return 1
    base_dir = args.submission_dir.resolve()
    text = md_path.read_text(encoding="utf-8")
    try:
        specs = parse_slides_md(text, base_dir)
    except SystemExit as e:
        logger.error("%s", e)
        return 1
    out = args.out.resolve()
    build_pptx(specs, out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
