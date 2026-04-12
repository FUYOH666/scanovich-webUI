"""Rich document ingest via markitdown: DOCX, XLSX, PPTX routing and conversion."""

from __future__ import annotations

import io
from typing import Any

import pytest

from gpthub_orchestrator.ingest.richdoc import (
    RichDocConvertError,
    convert_richdoc_bytes,
    is_richdoc_item,
)


# ---------------------------------------------------------------------------
# Detection
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "mime,filename",
    [
        ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", "report.docx"),
        ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "data.xlsx"),
        ("application/vnd.openxmlformats-officedocument.presentationml.presentation", "slides.pptx"),
        ("application/msword", "old.doc"),
        ("application/vnd.ms-excel", "old.xls"),
        ("application/rtf", "notes.rtf"),
        ("application/octet-stream", "report.docx"),  # wrong MIME but right extension
        ("application/octet-stream", "data.xlsx"),
        ("application/octet-stream", "slides.pptx"),
    ],
)
def test_is_richdoc_positive(mime: str, filename: str) -> None:
    assert is_richdoc_item(mime, filename)


@pytest.mark.parametrize(
    "mime,filename",
    [
        ("application/pdf", "report.pdf"),
        ("text/plain", "notes.txt"),
        ("text/csv", "data.csv"),
        ("application/json", "config.json"),
        ("audio/wav", "speech.wav"),
        ("image/png", "photo.png"),
        ("application/octet-stream", "unknown.bin"),
    ],
)
def test_is_richdoc_negative(mime: str, filename: str) -> None:
    assert not is_richdoc_item(mime, filename)


# ---------------------------------------------------------------------------
# Conversion — real DOCX built with python-docx
# ---------------------------------------------------------------------------


def _make_docx_bytes(text: str) -> bytes:
    """Create a minimal .docx in memory using python-docx."""
    try:
        from docx import Document  # type: ignore[import-not-found]
    except ImportError:
        pytest.skip("python-docx not installed")
    doc = Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx_bytes(data: list[list[Any]]) -> bytes:
    """Create a minimal .xlsx in memory using openpyxl."""
    try:
        from openpyxl import Workbook  # type: ignore[import-not-found]
    except ImportError:
        pytest.skip("openpyxl not installed")
    wb = Workbook()
    ws = wb.active
    assert ws is not None
    for row in data:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx_bytes(title: str, body: str) -> bytes:
    """Create a minimal .pptx in memory using python-pptx."""
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = body
            break
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_convert_docx_extracts_text() -> None:
    data = _make_docx_bytes("Привет, это тестовый документ для GPTHub!")
    result = convert_richdoc_bytes(data, filename="test.docx")
    assert "тестовый документ" in result


def test_convert_xlsx_extracts_cells() -> None:
    data = _make_xlsx_bytes([
        ["Name", "Score"],
        ["Alice", 95],
        ["Bob", 87],
    ])
    result = convert_richdoc_bytes(data, filename="scores.xlsx")
    assert "Alice" in result
    assert "Bob" in result


def test_convert_pptx_extracts_slide_text() -> None:
    data = _make_pptx_bytes("Architecture Overview", "Microservices and event sourcing")
    result = convert_richdoc_bytes(data, filename="deck.pptx")
    assert "Architecture" in result or "Microservices" in result


def test_convert_garbage_returns_fallback_text() -> None:
    """markitdown falls back to raw text for unrecognized content — not an error."""
    result = convert_richdoc_bytes(b"this is not a real document", filename="fake.docx")
    assert "not a real document" in result


def test_convert_empty_raises() -> None:
    with pytest.raises(RichDocConvertError):
        convert_richdoc_bytes(b"", filename="empty.docx")


# ---------------------------------------------------------------------------
# Pipeline integration — verify richdoc items get routed correctly
# ---------------------------------------------------------------------------


def test_pipeline_routes_docx_to_richdoc() -> None:
    """Verify that the pipeline routing logic picks richdoc for .docx files."""
    from gpthub_orchestrator.ingest.pipeline import _is_plain_text_item  # type: ignore[reportPrivateUsage]

    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    fn = "report.docx"
    # richdoc should match, but plain_text should NOT match (so we don't double-process)
    assert is_richdoc_item(mime, fn)
    assert not _is_plain_text_item(mime, fn)


def test_pipeline_routes_xlsx_to_richdoc() -> None:
    from gpthub_orchestrator.ingest.pipeline import _is_plain_text_item  # type: ignore[reportPrivateUsage]

    mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    fn = "data.xlsx"
    assert is_richdoc_item(mime, fn)
    assert not _is_plain_text_item(mime, fn)
