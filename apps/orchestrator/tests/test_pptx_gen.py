"""WOW-3 PPTX generation: intent, plan parser, builder, end-to-end."""

from __future__ import annotations

import json
import zipfile
from pathlib import Path
from typing import Any

import httpx
import pytest

from gpthub_orchestrator.pptx_gen import (
    DeckPlan,
    PptxPlanError,
    PptxResult,
    SlidePlan,
    _strip_code_fence,
    _strip_cot_blocks,
    _try_parse_plan_json,  # type: ignore[reportPrivateUsage]
    build_pptx_chat_completion,
    build_pptx_from_plan,
    build_pptx_message_text,
    build_pptx_sse_chunks,
    extract_pptx_topic,
    is_pptx_request,
    is_safe_token,
    public_pptx_url,
    request_slide_plan,
    resolve_pptx_path,
    run_pptx_generation,
    save_pptx_bytes,
    validate_plan,
)
from gpthub_orchestrator.settings import Settings


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------


def _mk_settings(tmp_path: Path | None = None, **over: Any) -> Settings:
    base: dict[str, Any] = {
        "litellm_base_url": "http://litellm:4000",
        "orchestrator_api_key": "t",
        "mws_gpt_api_base": "https://api.gpt.mws.ru/v1",
        "mws_gpt_api_key": "sk-test",
        "pptx_plan_timeout_seconds": 10.0,
        "pptx_min_slides": 3,
        "pptx_max_slides": 8,
        "pptx_public_base_url": "http://localhost:8089",
    }
    if tmp_path is not None:
        base["pptx_storage_dir"] = str(tmp_path / "pptx")
    base.update(over)
    return Settings(**base)  # type: ignore[arg-type]


def _chat_response(content: str) -> dict[str, Any]:
    return {
        "id": "chatcmpl-mock",
        "object": "chat.completion",
        "created": 0,
        "model": "mock",
        "choices": [
            {
                "index": 0,
                "message": {"role": "assistant", "content": content},
                "finish_reason": "stop",
            }
        ],
        "usage": {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0},
    }


_HAPPY_PLAN_JSON = json.dumps(
    {
        "title": "RAG в продакшене",
        "subtitle": "Практический гид",
        "slides": [
            {
                "title": "Введение",
                "bullets": ["Что такое RAG", "Зачем он бизнесу", "Где работает плохо"],
            },
            {
                "title": "Архитектура",
                "bullets": ["Indexer", "Retriever", "Re-ranker", "LLM"],
            },
            {
                "title": "Метрики качества",
                "bullets": ["Recall@k", "Faithfulness", "Latency", "Cost / token"],
            },
            {
                "title": "Подводные камни",
                "bullets": ["Сегментация", "Стейл данные", "Hallucinations"],
            },
            {
                "title": "Выводы",
                "bullets": ["Замеряй", "Оценивай", "Версионируй"],
            },
        ],
    },
    ensure_ascii=False,
)


# ---------------------------------------------------------------------------
# Intent detection
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "text",
    [
        "/pptx RAG в продакшене",
        "/slides MLOps стек",
        "Сделай презентацию по архитектуре нашей системы",
        "Создай презентацию о трендах LLM 2026",
        "Подготовь презентацию на тему feature store",
        "Презентация про vector databases",
        "Слайды по нашему MVP",
        "make a presentation about transformer scaling",
        "build a deck on enterprise RAG",
        "create slides for our roadmap",
        "powerpoint about AI safety",
        "давай, в формате pptx будет? создай",
        "сделай в формате pptx",
        "в pptx формате собери",
    ],
)
def test_is_pptx_request_positive(text: str) -> None:
    assert is_pptx_request(text)


@pytest.mark.parametrize(
    "text",
    [
        "Как работает RAG в двух предложениях",
        "Напиши код для feature store на Python",
        "Нарисуй схему архитектуры",
        "Привет!",
        "",
        "x" * 9000,
        "Расскажи про презентационные навыки",  # noun without verb context
    ],
)
def test_is_pptx_request_negative(text: str) -> None:
    assert not is_pptx_request(text)


def test_extract_pptx_topic_strips_slash() -> None:
    assert extract_pptx_topic("/pptx vector databases") == "vector databases"
    assert extract_pptx_topic("/slides MLOps стек") == "MLOps стек"


def test_extract_pptx_topic_passthrough() -> None:
    q = "Сделай презентацию по архитектуре нашей системы"
    assert extract_pptx_topic(q) == q


# ---------------------------------------------------------------------------
# JSON parsing helpers
# ---------------------------------------------------------------------------


def test_strip_code_fence_basic() -> None:
    raw = "```json\n{\"a\": 1}\n```"
    assert _strip_code_fence(raw) == '{"a": 1}'


def test_strip_code_fence_no_fence() -> None:
    assert _strip_code_fence('{"a": 1}') == '{"a": 1}'


def test_try_parse_plan_json_happy() -> None:
    obj = _try_parse_plan_json(_HAPPY_PLAN_JSON)
    assert isinstance(obj, dict)
    assert obj["title"] == "RAG в продакшене"


def test_try_parse_plan_json_with_fence() -> None:
    wrapped = f"```json\n{_HAPPY_PLAN_JSON}\n```"
    obj = _try_parse_plan_json(wrapped)
    assert obj["title"] == "RAG в продакшене"


def test_try_parse_plan_json_extracts_brace_block() -> None:
    noisy = "Sure, here you go:\n" + _HAPPY_PLAN_JSON + "\nLet me know if you need changes."
    obj = _try_parse_plan_json(noisy)
    assert obj["title"] == "RAG в продакшене"


def test_try_parse_plan_json_strips_think_block() -> None:
    """MWS glm-4.6 wraps output in <think>…</think> CoT before the actual JSON."""
    wrapped = "<think>Let me plan this...\n{fake braces}\n</think>\n" + _HAPPY_PLAN_JSON
    obj = _try_parse_plan_json(wrapped)
    assert obj["title"] == "RAG в продакшене"


def test_try_parse_plan_json_tolerates_trailing_prose() -> None:
    """Instruct models may append text after valid JSON (json.loads 'Extra data')."""
    noisy = _HAPPY_PLAN_JSON + "\n\nHope this helps! Summary follows…"
    obj = _try_parse_plan_json(noisy)
    assert isinstance(obj, dict)
    assert obj["title"] == "RAG в продакшене"


def test_try_parse_plan_json_fails_on_garbage() -> None:
    with pytest.raises(PptxPlanError):
        _try_parse_plan_json("definitely not json")


# ---------------------------------------------------------------------------
# Plan validation
# ---------------------------------------------------------------------------


def test_validate_plan_happy() -> None:
    raw = json.loads(_HAPPY_PLAN_JSON)
    plan = validate_plan(raw, min_slides=3, max_slides=10)
    assert plan.title == "RAG в продакшене"
    assert plan.subtitle == "Практический гид"
    assert len(plan.slides) == 5
    assert plan.slides[0].title == "Введение"
    assert plan.slides[0].bullets[0] == "Что такое RAG"


def test_validate_plan_caps_slides() -> None:
    raw = {"title": "T", "slides": [{"title": f"s{i}", "bullets": ["b"]} for i in range(20)]}
    plan = validate_plan(raw, min_slides=3, max_slides=8)
    assert len(plan.slides) == 8


def test_validate_plan_drops_empty_titles() -> None:
    raw = {
        "title": "T",
        "slides": [
            {"title": "ok", "bullets": ["a"]},
            {"title": "", "bullets": ["x"]},
            {"title": "second", "bullets": ["b"]},
            {"title": "third", "bullets": ["c"]},
        ],
    }
    plan = validate_plan(raw, min_slides=3, max_slides=8)
    assert [s.title for s in plan.slides] == ["ok", "second", "third"]


def test_validate_plan_missing_title_raises() -> None:
    with pytest.raises(PptxPlanError):
        validate_plan({"slides": [{"title": "a", "bullets": []}]}, min_slides=1, max_slides=5)


def test_validate_plan_missing_slides_raises() -> None:
    with pytest.raises(PptxPlanError):
        validate_plan({"title": "T"}, min_slides=1, max_slides=5)


def test_validate_plan_too_few_slides_raises() -> None:
    raw = {"title": "T", "slides": [{"title": "a", "bullets": ["b"]}]}
    with pytest.raises(PptxPlanError):
        validate_plan(raw, min_slides=3, max_slides=8)


def test_validate_plan_not_object_raises() -> None:
    with pytest.raises(PptxPlanError):
        validate_plan(["not", "a", "dict"], min_slides=1, max_slides=5)


def test_validate_plan_truncates_long_strings() -> None:
    long_title = "x" * 500
    raw = {
        "title": long_title,
        "slides": [
            {"title": long_title, "bullets": ["x" * 500]},
            {"title": "two", "bullets": ["b"]},
            {"title": "three", "bullets": ["c"]},
        ],
    }
    plan = validate_plan(raw, min_slides=3, max_slides=8)
    assert len(plan.title) <= 120
    assert plan.title.endswith("…")
    assert len(plan.slides[0].bullets[0]) <= 240


# ---------------------------------------------------------------------------
# python-pptx builder
# ---------------------------------------------------------------------------


def test_build_pptx_from_plan_returns_zip_bytes() -> None:
    plan = DeckPlan(
        title="Test deck",
        subtitle="A subtitle",
        slides=[
            SlidePlan(title="One", bullets=["a", "b", "c"]),
            SlidePlan(title="Two", bullets=["d"]),
            SlidePlan(title="Three", bullets=[]),
        ],
    )
    data = build_pptx_from_plan(plan)
    assert isinstance(data, bytes)
    assert len(data) > 1000
    # PPTX files are ZIP packages — verify the magic bytes.
    assert data[:4] == b"PK\x03\x04"
    # And that python-pptx round-trips it cleanly.
    import io

    with zipfile.ZipFile(io.BytesIO(data)) as z:
        names = z.namelist()
        assert "ppt/presentation.xml" in names
        assert any(n.startswith("ppt/slides/slide") for n in names)


def test_build_pptx_from_plan_round_trips_text() -> None:
    """Open the rendered .pptx via python-pptx and check titles match."""
    from pptx import Presentation  # type: ignore[import-not-found]

    plan = DeckPlan(
        title="Round Trip",
        subtitle="Sub",
        slides=[
            SlidePlan(title="Alpha", bullets=["b1", "b2"]),
            SlidePlan(title="Beta", bullets=["b3"]),
            SlidePlan(title="Gamma", bullets=[]),
        ],
    )
    data = build_pptx_from_plan(plan)
    import io

    prs = Presentation(io.BytesIO(data))
    # Title slide + 3 bullet slides = 4.
    assert len(prs.slides) == 4
    titles = []
    for s in prs.slides:
        if s.shapes.title is not None:
            titles.append(s.shapes.title.text)
    assert "Round Trip" in titles
    assert "Alpha" in titles
    assert "Beta" in titles
    assert "Gamma" in titles


# ---------------------------------------------------------------------------
# Storage
# ---------------------------------------------------------------------------


def test_save_and_resolve_pptx_bytes(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path)
    token, path = save_pptx_bytes(b"fake-pptx", settings=s)
    assert is_safe_token(token)
    assert Path(path).read_bytes() == b"fake-pptx"
    resolved = resolve_pptx_path(token, settings=s)
    assert resolved == path


def test_resolve_pptx_path_rejects_traversal(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path)
    assert resolve_pptx_path("../../etc/passwd", settings=s) is None
    assert resolve_pptx_path("../sneaky", settings=s) is None
    assert resolve_pptx_path("not-hex-token!", settings=s) is None


def test_resolve_pptx_path_returns_none_for_missing(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path)
    # 32 hex chars but no file on disk.
    assert resolve_pptx_path("a" * 32, settings=s) is None


def test_public_pptx_url_concat(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path, pptx_public_base_url="http://example.com/")
    assert public_pptx_url("abcd1234" * 4, settings=s) == "http://example.com/v1/files/pptx/" + "abcd1234" * 4


# ---------------------------------------------------------------------------
# request_slide_plan with mocked LiteLLM
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_request_slide_plan_happy(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path)
    calls: list[dict[str, Any]] = []

    def handler(request: httpx.Request) -> httpx.Response:
        body = json.loads(request.read())
        calls.append(body)
        return httpx.Response(200, json=_chat_response(_HAPPY_PLAN_JSON))

    transport = httpx.MockTransport(handler)
    async with httpx.AsyncClient(transport=transport) as http:
        plan = await request_slide_plan(
            http,
            settings=s,
            base_url="http://litellm:4000",
            auth_header="Bearer t",
            topic="RAG в продакшене",
        )
    assert plan.title == "RAG в продакшене"
    assert len(plan.slides) == 5
    assert len(calls) == 1
    assert calls[0]["model"] == s.pptx_plan_model


@pytest.mark.asyncio
async def test_request_slide_plan_retry_after_garbage(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path)
    answers = ["definitely not json", _HAPPY_PLAN_JSON]
    call_count = {"n": 0}

    def handler(request: httpx.Request) -> httpx.Response:
        idx = call_count["n"]
        call_count["n"] += 1
        return httpx.Response(200, json=_chat_response(answers[idx]))

    transport = httpx.MockTransport(handler)
    async with httpx.AsyncClient(transport=transport) as http:
        plan = await request_slide_plan(
            http,
            settings=s,
            base_url="http://litellm:4000",
            auth_header="",
            topic="anything",
        )
    assert call_count["n"] == 2
    assert plan.title == "RAG в продакшене"


@pytest.mark.asyncio
async def test_request_slide_plan_retry_then_fail(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path)

    def handler(request: httpx.Request) -> httpx.Response:
        return httpx.Response(200, json=_chat_response("still not json"))

    transport = httpx.MockTransport(handler)
    async with httpx.AsyncClient(transport=transport) as http:
        with pytest.raises(PptxPlanError):
            await request_slide_plan(
                http,
                settings=s,
                base_url="http://litellm:4000",
                auth_header="",
                topic="anything",
            )


@pytest.mark.asyncio
async def test_request_slide_plan_upstream_error(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path)

    def handler(request: httpx.Request) -> httpx.Response:
        return httpx.Response(502, json={"error": {"message": "boom"}})

    transport = httpx.MockTransport(handler)
    async with httpx.AsyncClient(transport=transport) as http:
        with pytest.raises(PptxPlanError):
            await request_slide_plan(
                http,
                settings=s,
                base_url="http://litellm:4000",
                auth_header="",
                topic="anything",
            )


@pytest.mark.asyncio
async def test_request_slide_plan_empty_topic_raises(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path)
    transport = httpx.MockTransport(lambda r: httpx.Response(200, json=_chat_response("{}")))
    async with httpx.AsyncClient(transport=transport) as http:
        with pytest.raises(PptxPlanError):
            await request_slide_plan(
                http,
                settings=s,
                base_url="http://litellm:4000",
                auth_header="",
                topic="   ",
            )


# ---------------------------------------------------------------------------
# Top-level run_pptx_generation end-to-end
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_run_pptx_generation_end_to_end(tmp_path: Path) -> None:
    s = _mk_settings(tmp_path)

    def handler(request: httpx.Request) -> httpx.Response:
        return httpx.Response(200, json=_chat_response(_HAPPY_PLAN_JSON))

    transport = httpx.MockTransport(handler)
    async with httpx.AsyncClient(transport=transport) as http:
        result = await run_pptx_generation(
            http,
            settings=s,
            base_url="http://litellm:4000",
            auth_header="Bearer t",
            topic="RAG в продакшене",
        )
    assert isinstance(result, PptxResult)
    assert result.plan.title == "RAG в продакшене"
    assert is_safe_token(result.token)
    assert result.url.endswith(result.token)
    assert result.byte_size > 1000
    # File actually written.
    on_disk = resolve_pptx_path(result.token, settings=s)
    assert on_disk is not None
    assert Path(on_disk).exists()
    # Trace payload shape.
    fb = result.trace_payload()
    assert fb["short_circuit"] == "pptx_generation"
    assert fb["title"] == "RAG в продакшене"
    assert fb["slide_count"] == 5
    assert fb["url"].endswith(result.token)


# ---------------------------------------------------------------------------
# Response builders + classifier wiring
# ---------------------------------------------------------------------------


def test_build_pptx_message_text_contains_link() -> None:
    plan = DeckPlan(
        title="Decko",
        subtitle="Sub",
        slides=[SlidePlan(title="One", bullets=["a"]), SlidePlan(title="Two", bullets=[])],
    )
    res = PptxResult(
        plan=plan,
        token="a" * 32,
        url="http://localhost:8089/v1/files/pptx/" + "a" * 32,
        byte_size=42_000,
        plan_ms=1000,
        build_ms=20,
    )
    text = build_pptx_message_text(res)
    assert "Decko" in text
    assert "Sub" in text
    assert "[Скачать .pptx" in text
    assert "/v1/files/pptx/" in text
    assert "1. One" in text
    assert "2. Two" in text


def test_build_pptx_chat_completion_shape() -> None:
    out = build_pptx_chat_completion(model_label="gpt-hub", text="hello")
    assert out["object"] == "chat.completion"
    assert out["model"] == "gpt-hub"
    assert out["choices"][0]["message"]["content"] == "hello"
    assert out["choices"][0]["finish_reason"] == "stop"
    assert out["id"].startswith("chatcmpl-pptx-")


def test_build_pptx_sse_chunks_terminates_with_done() -> None:
    chunks = build_pptx_sse_chunks("gpt-hub", "hello")
    assert len(chunks) == 3
    assert chunks[-1] == b"data: [DONE]\n\n"
    first = json.loads(chunks[0].removeprefix(b"data: ").strip())
    assert first["choices"][0]["delta"]["content"] == "hello"


def test_classifier_routes_pptx_intent_to_pptx_task() -> None:
    from gpthub_orchestrator.classifier import classify_messages

    msgs = [{"role": "user", "content": "Сделай презентацию по архитектуре нашей системы"}]
    out = classify_messages(msgs)
    assert out["task_type"] == "pptx_generation"


def test_classifier_normal_message_is_not_pptx() -> None:
    from gpthub_orchestrator.classifier import classify_messages

    msgs = [{"role": "user", "content": "Расскажи про vector databases"}]
    out = classify_messages(msgs)
    assert out["task_type"] != "pptx_generation"


def test_router_pptx_fallback_when_council_disabled() -> None:
    from gpthub_orchestrator.router import choose_model

    s = _mk_settings()
    classification = {"modalities": ["text"], "task_type": "pptx_generation"}
    out = choose_model(classification, s)
    assert out["reason"] == "pptx_generation_fallback"
