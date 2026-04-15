"""Tests for ``pptx.build`` — deck bytes from ``SlidePlan``."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation  # type: ignore[import-untyped]

from gpthub_orchestrator.pptx import PptxGenError, build_pptx_from_plan
from gpthub_orchestrator.pptx.build import _pick_template_path, deck_title_for_intro
from gpthub_orchestrator.pptx.schema import SlidePlan, SlideSpec
from gpthub_orchestrator.settings import Settings


def test_build_pptx_from_plan_zip_magic(pptx_settings: Settings) -> None:
    plan = SlidePlan(
        slides=[
            SlideSpec(title="T", bullets=["one"], notes=""),
        ],
    )
    blob = build_pptx_from_plan(plan, settings=pptx_settings)
    assert blob.startswith(b"PK")
    assert len(Presentation(BytesIO(blob)).slides) == 2  # intro + content

    blob_no_intro = build_pptx_from_plan(
        plan,
        settings=pptx_settings.model_copy(update={"pptx_intro_slide_enabled": False}),
    )
    assert len(Presentation(BytesIO(blob_no_intro)).slides) == 1


def test_build_pptx_on_slide_progress_callback(pptx_settings: Settings) -> None:
    plan = SlidePlan(
        slides=[
            SlideSpec(title="A", bullets=["1"], notes=""),
            SlideSpec(title="B", bullets=["2"], notes=""),
        ],
    )
    seen: list[tuple[int, int]] = []

    def on_progress(cur: int, tot: int) -> None:
        seen.append((cur, tot))

    build_pptx_from_plan(plan, settings=pptx_settings, on_slide_progress=on_progress)
    # intro + 2 content slides => 3 callbacks, total always 3
    assert seen == [(1, 3), (2, 3), (3, 3)]


def test_deck_title_for_intro_prefers_presentation_title() -> None:
    plan = SlidePlan(
        presentation_title="Обзор GPTHub",
        slides=[SlideSpec(title="Проблема", bullets=["x"], notes="")],
    )
    assert deck_title_for_intro(plan) == "Обзор GPTHub"


def test_deck_title_for_intro_falls_back_to_first_slide() -> None:
    plan = SlidePlan(
        slides=[SlideSpec(title="Only section", bullets=[], notes="")],
    )
    assert deck_title_for_intro(plan) == "Only section"


def test_build_pptx_empty_plan_raises(pptx_settings: Settings) -> None:
    with pytest.raises(PptxGenError, match="empty_plan"):
        build_pptx_from_plan(SlidePlan(slides=[]), settings=pptx_settings)


def test_build_pptx_with_bundled_slidesgo_template() -> None:
    root = Path(__file__).resolve().parents[2]
    tdir = root / "assets" / "pttx"
    if not tdir.is_dir() or not any(tdir.glob("*.pptx")):
        pytest.skip("assets/pttx templates not present")
    plan = SlidePlan(
        slides=[
            SlideSpec(title="Slide one", bullets=["Point A"], notes=""),
            SlideSpec(title="Slide two", bullets=[], notes=""),
        ],
    )
    s = Settings(
        litellm_base_url="http://litellm.test",
        orchestrator_api_key="k",
        pptx_asset_templates_enabled=True,
        pptx_templates_dir=str(tdir),
        pptx_template_index=0,
        pptx_plan_audience="general",
    )
    blob = build_pptx_from_plan(plan, settings=s)
    assert blob.startswith(b"PK")
    assert len(blob) > 500_000


def test_pick_template_path_by_audience() -> None:
    root = Path(__file__).resolve().parents[2]
    tdir = root / "assets" / "pttx"
    if not tdir.is_dir() or not (tdir / "dark.pptx").is_file():
        pytest.skip("dark.pptx not in assets/pttx")
    s = Settings(
        litellm_base_url="http://litellm.test",
        orchestrator_api_key="k",
        pptx_asset_templates_enabled=True,
        pptx_templates_dir=str(tdir),
        pptx_plan_audience="general",
        pptx_template_index=63,
    )
    picked = _pick_template_path(s)
    assert picked is not None
    assert picked.name == "dark.pptx"


def test_pick_template_path_fallback_index_when_file_missing(tmp_path: Path) -> None:
    only = tmp_path / "only.pptx"
    only.write_bytes(b"not a real pptx")
    s = Settings(
        litellm_base_url="http://litellm.test",
        orchestrator_api_key="k",
        pptx_asset_templates_enabled=True,
        pptx_templates_dir=str(tmp_path),
        pptx_plan_audience="general",
        pptx_template_index=0,
    )
    picked = _pick_template_path(s)
    assert picked is not None
    assert picked.name == "only.pptx"
