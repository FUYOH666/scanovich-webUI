"""Live LLM benchmark: deck «природа» against every ``assets/pttx`` template.

Opt-in only (cost + latency). Run from ``apps/orchestrator``::

    PPTX_BENCH=1 LITELLM_BASE_URL=http://127.0.0.1:4000 ORCHESTRATOR_API_KEY=... \\
      uv run pytest tests/pptx_tests/test_bench_nature_templates.py -s

Without ``PPTX_BENCH=1`` the test is skipped so CI stays fast.

Generated files are written under ``scanovich-webUI/tmp/`` (gitignored), one ``.pptx`` per template.
"""

from __future__ import annotations

import json
import os
import re
import shutil
import time
from io import BytesIO
from pathlib import Path

import httpx
import pytest
from pptx import Presentation  # type: ignore[import-untyped]

from gpthub_orchestrator.pptx import build_pptx_from_plan, request_slide_plan
from gpthub_orchestrator.settings import Settings

_ORCH_ROOT = Path(__file__).resolve().parents[2]
_SCANOVICH_ROOT = Path(__file__).resolve().parents[4]
_TEMPLATES_DIR = _ORCH_ROOT / "assets" / "pttx"
_TMP_OUT_DIR = _SCANOVICH_ROOT / "tmp"


def _safe_filename_stem(name: str, *, max_len: int = 96) -> str:
    base = Path(name).stem
    s = re.sub(r"[^\w\-]+", "_", base, flags=re.UNICODE).strip("_")
    return (s or "template")[:max_len]

_NATURE_USER_RU = (
    "Сделай презентацию на тему природа: экосистемы, охрана природы, "
    "климат и взаимосвязи живых организмов. Язык ответа — русский."
)


def _bench_skip_if_disabled() -> None:
    if os.environ.get("PPTX_BENCH", "").strip() != "1":
        pytest.skip("Set PPTX_BENCH=1 to run live PPTX benchmark (see module docstring)")
    base = (os.environ.get("LITELLM_BASE_URL") or "").strip()
    key = (os.environ.get("ORCHESTRATOR_API_KEY") or "").strip()
    if not base:
        pytest.skip("PPTX_BENCH requires LITELLM_BASE_URL")
    if not key:
        pytest.skip("PPTX_BENCH requires ORCHESTRATOR_API_KEY")


def _bench_settings(*, templates_dir: Path, template_index: int) -> Settings:
    return Settings(
        litellm_base_url=os.environ["LITELLM_BASE_URL"].rstrip("/"),
        orchestrator_api_key=os.environ["ORCHESTRATOR_API_KEY"],
        pptx_asset_templates_enabled=True,
        pptx_templates_dir=str(templates_dir),
        pptx_template_index=template_index,
        # Bound wall time / token use for a repeatable bench; raise in env if needed.
        pptx_max_slides=6,
        litellm_timeout_seconds=float(os.environ.get("PPTX_BENCH_LITELLM_TIMEOUT", "600")),
    )


@pytest.mark.bench
@pytest.mark.asyncio
async def test_bench_nature_deck_all_pttx_templates(tmp_path_factory: pytest.TempPathFactory) -> None:
    """Generate a nature deck once per ``*.pptx`` in ``assets/pttx``; assert valid bytes + slide count."""
    _bench_skip_if_disabled()
    if not _TEMPLATES_DIR.is_dir():
        pytest.skip(f"missing templates dir: {_TEMPLATES_DIR}")
    template_files = sorted(_TEMPLATES_DIR.glob("*.pptx"))
    if not template_files:
        pytest.skip(f"no *.pptx under {_TEMPLATES_DIR}")

    _TMP_OUT_DIR.mkdir(parents=True, exist_ok=True)

    messages: list[dict[str, str]] = [{"role": "user", "content": _NATURE_USER_RU}]
    rows: list[dict[str, object]] = []
    timeout = httpx.Timeout(_bench_settings(templates_dir=_TEMPLATES_DIR, template_index=0).litellm_timeout_seconds)

    async with httpx.AsyncClient(timeout=timeout) as http:
        for idx, tpath in enumerate(template_files):
            # One template per dir so audience-based default (dark.pptx) does not force the same theme every time.
            one_dir = tmp_path_factory.mktemp(f"bench_pttx_{idx}")
            shutil.copy(tpath, one_dir / tpath.name)
            settings = _bench_settings(templates_dir=one_dir, template_index=0)
            t0 = time.perf_counter()
            plan = await request_slide_plan(http, settings, messages)
            t_plan = time.perf_counter()
            blob = build_pptx_from_plan(plan, settings=settings)
            t_end = time.perf_counter()
            prs = Presentation(BytesIO(blob))
            n_slides = len(prs.slides)
            out_name = f"pptx_bench_nature_{idx:02d}_{_safe_filename_stem(tpath.name)}.pptx"
            out_path = _TMP_OUT_DIR / out_name
            out_path.write_bytes(blob)
            rows.append(
                {
                    "template_file": tpath.name,
                    "template_index": idx,
                    "saved_path": str(out_path.relative_to(_SCANOVICH_ROOT)),
                    "presentation_title": plan.presentation_title,
                    "plan_slides": len(plan.slides),
                    "pptx_slides": n_slides,
                    "plan_ms": round((t_plan - t0) * 1000, 1),
                    "build_ms": round((t_end - t_plan) * 1000, 1),
                    "total_ms": round((t_end - t0) * 1000, 1),
                    "pptx_bytes": len(blob),
                }
            )
            assert blob.startswith(b"PK"), tpath.name
            assert len(plan.slides) >= 1, tpath.name
            assert n_slides >= 1, tpath.name

    print("\nPPTX_BENCH summary:\n" + json.dumps(rows, ensure_ascii=False, indent=2))
