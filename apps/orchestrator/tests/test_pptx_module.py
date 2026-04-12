import json
import os
from io import BytesIO
from pathlib import Path

import httpx
import pytest
from pptx import Presentation  # type: ignore[import-untyped]

os.environ.setdefault("LITELLM_BASE_URL", "http://127.0.0.1:9")
os.environ.setdefault("ORCHESTRATOR_API_KEY", "k")

from gpthub_orchestrator.model_registry import load_model_roles
from gpthub_orchestrator.pptx import (
    PptxGenError,
    SlidePlan,
    SlideSpec,
    build_pptx_from_plan,
    extract_json_object,
    markdown_preview_with_download_link,
    parse_slide_plan_text,
    request_slide_plan,
)
from gpthub_orchestrator.pptx.parse import (
    SLIDE_AGENT_MAX_VISIBLE_CHARS,
    clamp_slide_visible_to_max,
)
from gpthub_orchestrator.settings import Settings


@pytest.fixture(autouse=True)
def clear_registry_cache():
    load_model_roles.cache_clear()
    yield
    load_model_roles.cache_clear()


def _settings(**kwargs):
    base = {
        "litellm_base_url": "http://litellm.test",
        "orchestrator_api_key": "k",
        "pptx_asset_templates_enabled": False,
    }
    base.update(kwargs)
    return Settings(**base)


def test_extract_json_object_fence():
    raw = 'Sure\n```json\n{"slides":[]}\n```\n'
    assert '"slides"' in extract_json_object(raw)


def test_parse_slide_plan_text_minimal():
    text = json.dumps(
        {
            "slides": [
                {"title": "Intro", "bullets": ["a", "b"], "notes": "say hi"},
            ]
        }
    )
    plan = parse_slide_plan_text(text)
    assert len(plan.slides) == 1
    assert plan.slides[0].title == "Intro"
    assert plan.slides[0].bullets == ["a", "b"]
    assert plan.slides[0].kind is None


def test_parse_slide_plan_text_with_kind():
    text = json.dumps(
        {
            "slides": [
                {"title": "Roadmap", "bullets": ["Q1"], "notes": "", "kind": "timeline"},
            ]
        }
    )
    plan = parse_slide_plan_text(text)
    assert plan.slides[0].kind == "timeline"


def test_clamp_slide_visible_drops_bullets_from_end_until_under_cap():
    """Whole bullets removed from the tail until title + bullets ≤ max."""
    title = "T" * 50
    bullets = ["x" * 100, "y" * 100, "z" * 100, "w" * 100, "last" * 20]
    spec = SlideSpec(title=title, bullets=bullets, notes="keep", kind="bullets")
    out = clamp_slide_visible_to_max(spec, max_chars=SLIDE_AGENT_MAX_VISIBLE_CHARS)
    assert out.notes == "keep"
    assert out.kind == "bullets"
    total = len(out.title) + sum(len(b) for b in out.bullets)
    assert total <= SLIDE_AGENT_MAX_VISIBLE_CHARS
    assert len(out.bullets) < len(bullets)


def test_clamp_slide_visible_truncates_last_bullet_when_one_remains():
    """Single long line without spaces: last resort is a prefix cap (no mid-word cut when possible)."""
    title = "Hi"
    long_bullet = "B" * 600
    spec = SlideSpec(title=title, bullets=[long_bullet], notes="", kind=None)
    out = clamp_slide_visible_to_max(spec, max_chars=SLIDE_AGENT_MAX_VISIBLE_CHARS)
    assert len(out.title) + sum(len(b) for b in out.bullets) <= SLIDE_AGENT_MAX_VISIBLE_CHARS
    want = SLIDE_AGENT_MAX_VISIBLE_CHARS - len(title)
    assert out.bullets == [long_bullet[:want]]


def test_clamp_slide_visible_truncates_long_line_at_word_boundary():
    title = "X"
    words = "word " * 300  # > 500 chars total with title
    spec = SlideSpec(title=title, bullets=[words], notes="", kind=None)
    out = clamp_slide_visible_to_max(spec)
    b = out.bullets[0]
    assert len(title) + len(b) <= SLIDE_AGENT_MAX_VISIBLE_CHARS
    assert b.strip()
    assert not b.endswith("wo")
    assert b.endswith(" ") or b.endswith("word")


def test_clamp_slide_visible_drops_trailing_lines_in_bullet():
    """Newlines separate lines; trailing whole lines are removed before mid-line cuts."""
    title = "T"
    bullet = "first line short\n" + ("second " * 120) + "\n" + "third should drop entirely"
    spec = SlideSpec(title=title, bullets=[bullet], notes="", kind=None)
    out = clamp_slide_visible_to_max(spec)
    assert "third should drop" not in out.bullets[0]
    assert len(title) + len(out.bullets[0]) <= SLIDE_AGENT_MAX_VISIBLE_CHARS


def test_clamp_slide_visible_truncates_title_when_no_bullets():
    long_title = "N" * 600
    spec = SlideSpec(title=long_title, bullets=[], notes="", kind=None)
    out = clamp_slide_visible_to_max(spec)
    assert out.title == long_title[:SLIDE_AGENT_MAX_VISIBLE_CHARS]
    assert out.bullets == []


def test_clamp_slide_visible_title_priority_when_budget_exhausted():
    """If title alone exceeds remaining budget with one bullet, title wins (bullet dropped)."""
    title = "T" * 520
    spec = SlideSpec(title=title, bullets=["extra"], notes="", kind=None)
    out = clamp_slide_visible_to_max(spec)
    assert out.title == title[:SLIDE_AGENT_MAX_VISIBLE_CHARS]
    assert out.bullets == []


def test_clamp_slide_visible_unchanged_when_already_short():
    spec = SlideSpec(title="A", bullets=["b", "c"], notes="n", kind="stats")
    out = clamp_slide_visible_to_max(spec)
    assert out == spec


def test_parse_slide_plan_text_invalid_kind_dropped():
    text = json.dumps(
        {
            "slides": [
                {"title": "X", "bullets": [], "notes": "", "kind": "not-a-real-layout"},
            ]
        }
    )
    plan = parse_slide_plan_text(text)
    assert plan.slides[0].kind is None


def test_markdown_preview_shows_kind():
    plan = SlidePlan(
        slides=[SlideSpec(title="T", bullets=["a"], notes="", kind="bullets")],
    )
    md = markdown_preview_with_download_link(plan, "https://example.test/d.pptx?token=x")
    assert "макет: `bullets`" in md


def test_markdown_preview_intro_line():
    plan = SlidePlan(
        slides=[SlideSpec(title="Тема", bullets=["a"], notes="")],
    )
    md = markdown_preview_with_download_link(
        plan, "https://example.test/d.pptx?token=x", intro_title="Тема"
    )
    assert "Титульный слайд" in md
    assert "Тема" in md


def test_build_pptx_from_plan_zip_magic():
    plan = SlidePlan(
        slides=[
            SlideSpec(title="T", bullets=["one"], notes=""),
        ],
    )
    blob = build_pptx_from_plan(plan, settings=_settings())
    assert blob.startswith(b"PK")
    assert len(Presentation(BytesIO(blob)).slides) == 2  # intro + content

    blob_no_intro = build_pptx_from_plan(
        plan, settings=_settings(pptx_intro_slide_enabled=False)
    )
    assert len(Presentation(BytesIO(blob_no_intro)).slides) == 1


def test_build_pptx_empty_plan_raises():
    with pytest.raises(PptxGenError, match="empty_plan"):
        build_pptx_from_plan(SlidePlan(slides=[]), settings=_settings())


@pytest.mark.asyncio
async def test_request_slide_plan_success():
    good = json.dumps(
        {"slides": [{"title": "S1", "bullets": ["x"], "notes": ""}]}
    )

    def handler(request: httpx.Request) -> httpx.Response:
        assert request.url.path.endswith("/v1/chat/completions")
        body = json.loads(request.content.decode())
        sys0 = body["messages"][0]["content"]
        assert "Tone: auto" in sys0
        assert "Text density level" in sys0
        assert "timeline" in sys0
        return httpx.Response(
            200,
            json={"choices": [{"message": {"content": good}}]},
        )

    transport = httpx.MockTransport(handler)
    async with httpx.AsyncClient(transport=transport) as http:
        plan = await request_slide_plan(
            http,
            _settings(),
            [{"role": "user", "content": "Сделай презентацию про тесты"}],
        )
    assert plan.slides[0].title == "S1"


@pytest.mark.asyncio
async def test_request_slide_plan_json_retry_second_turn():
    calls = {"n": 0}
    good = json.dumps(
        {"slides": [{"title": "Fixed", "bullets": [], "notes": ""}]}
    )

    def handler(request: httpx.Request) -> httpx.Response:
        calls["n"] += 1
        if calls["n"] == 1:
            return httpx.Response(
                200,
                json={"choices": [{"message": {"content": "not json at all"}}]},
            )
        return httpx.Response(
            200,
            json={"choices": [{"message": {"content": good}}]},
        )

    transport = httpx.MockTransport(handler)
    async with httpx.AsyncClient(transport=transport) as http:
        plan = await request_slide_plan(
            http,
            _settings(),
            [{"role": "user", "content": "/pptx about QA"}],
        )
    assert calls["n"] == 2
    assert plan.slides[0].title == "Fixed"


@pytest.mark.asyncio
async def test_request_slide_plan_chain_429_then_ok():
    models: list[str] = []

    def handler(request: httpx.Request) -> httpx.Response:
        body = json.loads(request.content.decode())
        models.append(str(body.get("model", "")))
        if body.get("model") == "gpt-hub-strong":
            return httpx.Response(429, json={"error": "rate"})
        good = json.dumps({"slides": [{"title": "R", "bullets": [], "notes": ""}]})
        return httpx.Response(
            200,
            json={"choices": [{"message": {"content": good}}]},
        )

    transport = httpx.MockTransport(handler)
    async with httpx.AsyncClient(transport=transport) as http:
        plan = await request_slide_plan(
            http,
            _settings(),
            [{"role": "user", "content": "build deck for sprint review"}],
        )
    assert models[0] == "gpt-hub-strong"
    assert models[1] == "gpt-hub-turbo"
    assert plan.slides[0].title == "R"


def test_build_pptx_with_bundled_slidesgo_template():
    root = Path(__file__).resolve().parents[1]
    tdir = root / "assets" / "pttx"
    if not tdir.is_dir() or not any(tdir.glob("*.pptx")):
        pytest.skip("assets/pttx templates not present")
    plan = SlidePlan(
        slides=[
            SlideSpec(title="Slide one", bullets=["Point A"], notes=""),
            SlideSpec(title="Slide two", bullets=[], notes=""),
        ],
    )
    s = Settings(
        litellm_base_url="http://litellm.test",
        orchestrator_api_key="k",
        pptx_asset_templates_enabled=True,
        pptx_templates_dir=str(tdir),
        pptx_template_index=0,
    )
    blob = build_pptx_from_plan(plan, settings=s)
    assert blob.startswith(b"PK")
    assert len(blob) > 500_000
