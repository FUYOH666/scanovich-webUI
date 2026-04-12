"""PPTX generation short-circuit (WOW-3): topic → slide plan → .pptx file.

Design:
* Detects «сделай презентацию / build a deck / /pptx» on the last user
  text and short-circuits the chat path.
* Asks ``pptx_plan_model`` (default ``gpt-hub-strong`` → glm-4.6) to
  produce a strict JSON slide plan with one retry on parse failure.
* Builds the .pptx with python-pptx (title slide + bullet slides).
* Writes bytes to ``settings.pptx_storage_dir`` and returns one
  OpenAI-compatible ``chat.completion`` containing a markdown link to a
  GET endpoint exposed by ``main.py`` (no auth, unguessable URL).
* Stays inside the «one chat → one answer» contract: no second turn,
  no streaming-only special case, no leaking JSON to the user.

The plan JSON contract is intentionally tiny so the model rarely fails
validation::

    {
      "title": "Presentation title",
      "subtitle": "Optional one-liner",
      "slides": [
        {"title": "Slide title", "bullets": ["bullet a", "bullet b"]}
      ]
    }
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import time
import uuid
from dataclasses import dataclass, field
from typing import Any

import httpx

from gpthub_orchestrator.reasoning_response_filter import merge_reasoning_exclude_into_body
from gpthub_orchestrator.settings import Settings

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Intent detection
# ---------------------------------------------------------------------------

# Slash commands always win.
_SLASH_PPTX = re.compile(r"(?:^|\s)/(?:pptx|slides|deck|presentation)\b", re.IGNORECASE)

_PPTX_PHRASES = [
    # RU verbs that almost always mean "make a deck".
    re.compile(
        r"\b(?:сделай|создай|сгенерир\w*|подготов\w*|собер[иь]|построй|нарису\w*)"
        r"[^.?!\n]{0,40}\bпрезентац\w*",
        re.IGNORECASE | re.UNICODE,
    ),
    re.compile(r"\bпрезентац\w*\b[^.?!\n]{0,30}\b(?:по|про|о|об|на\s+тему)\b", re.IGNORECASE | re.UNICODE),
    re.compile(r"\bслайд\w*\s+(?:по|про|о|об|на\s+тему)\b", re.IGNORECASE | re.UNICODE),
    re.compile(r"\bколод\w*\s+слайд\w*\b", re.IGNORECASE | re.UNICODE),
    # EN.
    re.compile(
        r"\b(?:make|build|create|generate|draft|prepare)\s+"
        r"(?:an?\s+|the\s+)?(?:presentation|deck|slides|slide\s+deck|powerpoint|pptx)\b",
        re.IGNORECASE,
    ),
    re.compile(r"\bslide\s+deck\s+(?:on|about|for)\b", re.IGNORECASE),
    re.compile(r"\b(?:powerpoint|pptx)\s+(?:on|about|for)\b", re.IGNORECASE),
]


def is_pptx_request(text: str) -> bool:
    """True if the last user text should short-circuit into PPTX generation."""
    if not text:
        return False
    s = text.strip()
    if not s:
        return False
    # Hard ceiling: don't hijack a giant document paste into PPTX.
    if len(s) > 8000:
        return False
    if _SLASH_PPTX.search(s):
        return True
    for pat in _PPTX_PHRASES:
        if pat.search(s):
            return True
    return False


def extract_pptx_topic(text: str) -> str:
    """Strip a leading slash command. Return the original text otherwise.

    The plan model receives the full request, so we don't try to be clever
    about extracting just the topic — the model is good at that.
    """
    s = text.strip()
    m = _SLASH_PPTX.search(s)
    if m and m.start() <= 1:
        return s[m.end():].strip(" :—-")
    return s


# ---------------------------------------------------------------------------
# Plan model + validation
# ---------------------------------------------------------------------------


@dataclass
class SlidePlan:
    title: str
    bullets: list[str] = field(default_factory=list)


@dataclass
class DeckPlan:
    title: str
    subtitle: str
    slides: list[SlidePlan] = field(default_factory=list)


class PptxPlanError(ValueError):
    """Raised when the plan model output cannot be parsed into a DeckPlan."""


def _coerce_str(v: Any, *, max_len: int) -> str:
    if v is None:
        return ""
    s = str(v).strip()
    if len(s) > max_len:
        s = s[: max_len - 1].rstrip() + "…"
    return s


def validate_plan(raw: Any, *, min_slides: int, max_slides: int) -> DeckPlan:
    """Strict-but-lenient validator: required keys present, lists trimmed."""
    if not isinstance(raw, dict):
        raise PptxPlanError("plan_not_object")
    title = _coerce_str(raw.get("title"), max_len=120)
    if not title:
        raise PptxPlanError("plan_missing_title")
    subtitle = _coerce_str(raw.get("subtitle"), max_len=200)
    slides_raw = raw.get("slides")
    if not isinstance(slides_raw, list) or not slides_raw:
        raise PptxPlanError("plan_missing_slides")
    slides: list[SlidePlan] = []
    for item in slides_raw:
        if not isinstance(item, dict):
            continue
        s_title = _coerce_str(item.get("title"), max_len=120)
        if not s_title:
            continue
        bullets_raw = item.get("bullets") or item.get("points") or []
        if not isinstance(bullets_raw, list):
            bullets_raw = []
        bullets: list[str] = []
        for b in bullets_raw[:8]:
            bs = _coerce_str(b, max_len=240)
            if bs:
                bullets.append(bs)
        slides.append(SlidePlan(title=s_title, bullets=bullets))
        if len(slides) >= max_slides:
            break
    if len(slides) < min_slides:
        raise PptxPlanError(f"plan_too_few_slides: got {len(slides)}, need {min_slides}")
    return DeckPlan(title=title, subtitle=subtitle, slides=slides)


def _build_plan_system_prompt(*, min_slides: int, max_slides: int) -> str:
    return (
        "Ты — генератор плана презентации. Твоя единственная задача — вернуть "
        "СТРОГО ВАЛИДНЫЙ JSON-объект без какого-либо текста до или после, без "
        "```json ... ``` обёрток, без комментариев, без markdown.\n\n"
        "Контракт JSON:\n"
        '{"title": "Заголовок презентации (до 100 символов)", '
        '"subtitle": "Короткая строка-подзаголовок (опционально, до 180)", '
        '"slides": [{"title": "Заголовок слайда", "bullets": ["буллет 1", "буллет 2", "буллет 3"]}]}\n\n'
        "Правила:\n"
        f"- Минимум {min_slides} и максимум {max_slides} слайдов в массиве slides.\n"
        "- На каждом слайде 3–6 буллетов, каждый буллет — короткое предложение или фраза до 220 символов.\n"
        "- Язык — тот же, что и у запроса пользователя.\n"
        "- Не вставляй заголовок 'Содержание' / 'Спасибо' автоматически — только содержательные слайды.\n"
        "- Никаких эмодзи, никаких HTML-тегов, никакого markdown внутри значений.\n"
        "- Возвращай ТОЛЬКО JSON. Никаких пояснений."
    )

_PLAN_RETRY_PROMPT = (
    "Предыдущий ответ нельзя распарсить как JSON по контракту. Ошибка: {error}.\n"
    "Верни СТРОГО валидный JSON по контракту, без markdown и пояснений."
)


def _strip_code_fence(text: str) -> str:
    """Best-effort: peel ```json ... ``` fences if the model added them."""
    s = text.strip()
    if s.startswith("```"):
        # Drop the first fence line.
        nl = s.find("\n")
        if nl != -1:
            s = s[nl + 1 :]
        if s.endswith("```"):
            s = s[: -3]
    return s.strip()


def _try_parse_plan_json(text: str) -> Any:
    """Parse raw assistant content into a Python object.

    Tolerates ``` fences and a single trailing comma; raises
    PptxPlanError on hard failure.
    """
    s = _strip_code_fence(text)
    # Some models still wrap JSON in <json> ... </json> tags.
    s = re.sub(r"^<json>\s*", "", s, flags=re.IGNORECASE)
    s = re.sub(r"\s*</json>$", "", s, flags=re.IGNORECASE)
    try:
        return json.loads(s)
    except json.JSONDecodeError:
        # Last-resort: extract the largest {...} block.
        first = s.find("{")
        last = s.rfind("}")
        if first != -1 and last > first:
            try:
                return json.loads(s[first : last + 1])
            except json.JSONDecodeError as e2:
                raise PptxPlanError(f"json_parse_error: {e2}") from e2
        raise PptxPlanError("json_no_object")


async def _call_plan_chat(
    http: httpx.AsyncClient,
    *,
    base_url: str,
    auth_header: str,
    model: str,
    messages: list[dict[str, Any]],
    max_tokens: int,
    timeout_seconds: float,
) -> str:
    """Thin LiteLLM /v1/chat/completions wrapper for plan generation."""
    url = base_url.rstrip("/") + "/v1/chat/completions"
    payload: dict[str, Any] = {
        "model": model,
        "messages": messages,
        "max_tokens": max_tokens,
        "stream": False,
    }
    merge_reasoning_exclude_into_body(payload, enabled=True)
    headers: dict[str, str] = {"Content-Type": "application/json"}
    if auth_header:
        headers["Authorization"] = auth_header
    r = await http.post(
        url,
        headers=headers,
        json=payload,
        timeout=httpx.Timeout(timeout_seconds, connect=15.0),
    )
    if r.status_code >= 400:
        raise PptxPlanError(f"plan_status_{r.status_code}: {r.text[:200]}")
    try:
        data = r.json()
    except Exception as e:  # noqa: BLE001
        raise PptxPlanError("plan_invalid_json") from e
    choices = data.get("choices") if isinstance(data, dict) else None
    if not isinstance(choices, list) or not choices:
        raise PptxPlanError("plan_empty_choices")
    msg = choices[0].get("message") if isinstance(choices[0], dict) else None
    if not isinstance(msg, dict):
        raise PptxPlanError("plan_bad_message")
    content = msg.get("content")
    if not isinstance(content, str) or not content.strip():
        raise PptxPlanError("plan_empty_content")
    return content


async def request_slide_plan(
    http: httpx.AsyncClient,
    *,
    settings: Settings,
    base_url: str,
    auth_header: str,
    topic: str,
) -> DeckPlan:
    """Ask the plan model for a deck plan, with one retry on parse failure.

    Always raises :class:`PptxPlanError` on terminal failure so the caller
    can fall through to normal chat routing.
    """
    if not topic.strip():
        raise PptxPlanError("empty_topic")
    sys_prompt = _build_plan_system_prompt(
        min_slides=settings.pptx_min_slides,
        max_slides=settings.pptx_max_slides,
    )
    messages: list[dict[str, Any]] = [
        {"role": "system", "content": sys_prompt},
        {"role": "user", "content": topic.strip()},
    ]
    raw = await _call_plan_chat(
        http,
        base_url=base_url,
        auth_header=auth_header,
        model=settings.pptx_plan_model,
        messages=messages,
        max_tokens=settings.pptx_plan_max_tokens,
        timeout_seconds=settings.pptx_plan_timeout_seconds,
    )
    try:
        parsed = _try_parse_plan_json(raw)
        return validate_plan(
            parsed,
            min_slides=settings.pptx_min_slides,
            max_slides=settings.pptx_max_slides,
        )
    except PptxPlanError as e:
        logger.warning("pptx_plan_first_attempt_failed err=%s", e)
        retry_messages: list[dict[str, Any]] = [
            {"role": "system", "content": sys_prompt},
            {"role": "user", "content": topic.strip()},
            {"role": "assistant", "content": raw},
            {"role": "user", "content": _PLAN_RETRY_PROMPT.format(error=str(e))},
        ]
        raw2 = await _call_plan_chat(
            http,
            base_url=base_url,
            auth_header=auth_header,
            model=settings.pptx_plan_model,
            messages=retry_messages,
            max_tokens=settings.pptx_plan_max_tokens,
            timeout_seconds=settings.pptx_plan_timeout_seconds,
        )
        parsed2 = _try_parse_plan_json(raw2)
        return validate_plan(
            parsed2,
            min_slides=settings.pptx_min_slides,
            max_slides=settings.pptx_max_slides,
        )


# ---------------------------------------------------------------------------
# python-pptx builder
# ---------------------------------------------------------------------------


def build_pptx_from_plan(plan: DeckPlan) -> bytes:
    """Render a DeckPlan into raw .pptx bytes.

    Uses the default 16:9 template ships with python-pptx. Layout 0 is the
    title slide; layout 1 is the title + content layout we use for bullets.
    """
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation()
    # Title slide.
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    if title_slide.shapes.title is not None:
        title_slide.shapes.title.text = plan.title
    if plan.subtitle and len(title_slide.placeholders) > 1:
        try:
            title_slide.placeholders[1].text = plan.subtitle
        except (KeyError, IndexError):
            pass

    # Bullet slides.
    bullet_layout = prs.slide_layouts[1]
    for slide_plan in plan.slides:
        slide = prs.slides.add_slide(bullet_layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = slide_plan.title
        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_ph = ph
                break
        if body_ph is None:
            continue
        tf = body_ph.text_frame
        tf.clear()
        if not slide_plan.bullets:
            tf.text = ""
            continue
        # First bullet goes into the existing paragraph; subsequent bullets
        # are added as new paragraphs at level 0.
        first = slide_plan.bullets[0]
        tf.paragraphs[0].text = first
        tf.paragraphs[0].level = 0
        for b in slide_plan.bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Storage
# ---------------------------------------------------------------------------


_TOKEN_RE = re.compile(r"^[a-f0-9]{16,64}$")


def _safe_token() -> str:
    return uuid.uuid4().hex


def is_safe_token(token: str) -> bool:
    return bool(_TOKEN_RE.fullmatch(token or ""))


def save_pptx_bytes(data: bytes, *, settings: Settings) -> tuple[str, str]:
    """Persist .pptx bytes to disk; return ``(token, absolute_path)``."""
    os.makedirs(settings.pptx_storage_dir, exist_ok=True)
    token = _safe_token()
    path = os.path.join(settings.pptx_storage_dir, f"{token}.pptx")
    with open(path, "wb") as f:
        f.write(data)
    return token, path


def resolve_pptx_path(token: str, *, settings: Settings) -> str | None:
    """Resolve a token to an on-disk path. Returns None if invalid/missing."""
    if not is_safe_token(token):
        return None
    path = os.path.join(settings.pptx_storage_dir, f"{token}.pptx")
    if not os.path.isfile(path):
        return None
    return path


def public_pptx_url(token: str, *, settings: Settings) -> str:
    base = settings.pptx_public_base_url.rstrip("/")
    return f"{base}/v1/files/pptx/{token}"


# ---------------------------------------------------------------------------
# Result + OpenAI-compatible response builders
# ---------------------------------------------------------------------------


@dataclass
class PptxResult:
    plan: DeckPlan
    token: str
    url: str
    byte_size: int
    plan_ms: int
    build_ms: int

    def trace_payload(self) -> dict[str, Any]:
        return {
            "short_circuit": "pptx_generation",
            "plan_model": "",  # caller fills
            "title": self.plan.title,
            "slide_count": len(self.plan.slides),
            "byte_size": self.byte_size,
            "plan_ms": self.plan_ms,
            "build_ms": self.build_ms,
            "url": self.url,
        }


def build_pptx_message_text(result: PptxResult) -> str:
    lines: list[str] = []
    lines.append(f"**Готова презентация: {result.plan.title}**")
    if result.plan.subtitle:
        lines.append(f"_{result.plan.subtitle}_")
    lines.append("")
    lines.append(f"[Скачать .pptx ({result.byte_size // 1024} KB)]({result.url})")
    lines.append("")
    lines.append("**План презентации:**")
    for i, s in enumerate(result.plan.slides, start=1):
        lines.append(f"{i}. {s.title}")
    return "\n".join(lines)


def build_pptx_chat_completion(*, model_label: str, text: str) -> dict[str, Any]:
    return {
        "id": f"chatcmpl-pptx-{uuid.uuid4().hex[:12]}",
        "object": "chat.completion",
        "created": int(time.time()),
        "model": model_label,
        "choices": [
            {
                "index": 0,
                "message": {"role": "assistant", "content": text},
                "finish_reason": "stop",
            }
        ],
        "usage": {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0},
    }


def build_pptx_sse_chunks(model_label: str, text: str) -> list[bytes]:
    cid = f"chatcmpl-pptx-{uuid.uuid4().hex[:12]}"
    created = int(time.time())
    first = {
        "id": cid,
        "object": "chat.completion.chunk",
        "created": created,
        "model": model_label,
        "choices": [
            {
                "index": 0,
                "delta": {"role": "assistant", "content": text},
                "finish_reason": None,
            }
        ],
    }
    final = {
        "id": cid,
        "object": "chat.completion.chunk",
        "created": created,
        "model": model_label,
        "choices": [{"index": 0, "delta": {}, "finish_reason": "stop"}],
    }
    return [
        b"data: " + json.dumps(first, ensure_ascii=False).encode("utf-8") + b"\n\n",
        b"data: " + json.dumps(final, ensure_ascii=False).encode("utf-8") + b"\n\n",
        b"data: [DONE]\n\n",
    ]


# ---------------------------------------------------------------------------
# Top-level orchestration
# ---------------------------------------------------------------------------


async def run_pptx_generation(
    http: httpx.AsyncClient,
    *,
    settings: Settings,
    base_url: str,
    auth_header: str,
    topic: str,
) -> PptxResult:
    """Plan + build + save. Raises :class:`PptxPlanError` on plan failure.

    Storage / build are not retried on disk errors — they bubble up to the
    caller, which falls through to normal chat routing.
    """
    plan_t0 = time.monotonic()
    plan = await request_slide_plan(
        http,
        settings=settings,
        base_url=base_url,
        auth_header=auth_header,
        topic=topic,
    )
    plan_ms = int((time.monotonic() - plan_t0) * 1000)
    build_t0 = time.monotonic()
    data = build_pptx_from_plan(plan)
    build_ms = int((time.monotonic() - build_t0) * 1000)
    token, _path = save_pptx_bytes(data, settings=settings)
    return PptxResult(
        plan=plan,
        token=token,
        url=public_pptx_url(token, settings=settings),
        byte_size=len(data),
        plan_ms=plan_ms,
        build_ms=build_ms,
    )
